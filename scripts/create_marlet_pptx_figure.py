from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("figures/pptx/marlet_framework_figure.pptx")

W, H = 13.333, 7.5

INK = "172033"
MUTED = "64748B"
PAPER = "F7F9FC"
WHITE = "FFFFFF"
SLATE = "CBD5E1"
TEAL = "0F766E"
TEAL_LIGHT = "DDF7F4"
BLUE = "2563EB"
BLUE_LIGHT = "E8F0FF"
GREEN = "15803D"
GREEN_LIGHT = "E8F7ED"
GOLD = "B7791F"
GOLD_LIGHT = "FFF3D6"
RED = "B42318"
RED_LIGHT = "FFE9E6"
PURPLE = "6D28D9"
PURPLE_LIGHT = "F0E7FF"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_text_frame(shape, text: str, size: float, color: str = INK, bold: bool = False, align=PP_ALIGN.CENTER) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Aptos"
        run.font.size = Pt(size if i == 0 else max(size - 2, 6))
        run.font.bold = bold if i == 0 else False
        run.font.color.rgb = rgb(color if i == 0 else MUTED)


def rounded(slide, x, y, w, h, text, fill, line, size=9.5, bold=True, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(1.1)
    set_text_frame(shp, text, size=size, bold=bold)
    return shp


def label(slide, x, y, w, h, text, size=7.5, color=MUTED, bold=False):
    shp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text_frame(shp, text, size=size, color=color, bold=bold)
    return shp


def group_band(slide, x, y, w, h, fill, line, title):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = 20_000
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(0.8)
    label(slide, x + 0.12, y + 0.03, w - 0.24, 0.2, title, size=6.5, color=line, bold=True)
    return shp


def diamond(slide, x, y, w, h, text, fill=GOLD_LIGHT, line=GOLD, size=8.5):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(1.1)
    set_text_frame(shp, text, size=size, bold=True)
    return shp


def center(shape):
    return (shape.left + shape.width / 2, shape.top + shape.height / 2)


def add_line(slide, x1, y1, x2, y2, color=INK, width=1.2, dash=False, arrow=True):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    if dash:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        conn.line.end_arrowhead = True
    return conn


def add_polyline(slide, points, color=INK, width=1.2, dash=False):
    for (x1, y1), (x2, y2) in zip(points, points[1:]):
        add_line(slide, x1, y1, x2, y2, color=color, width=width, dash=dash, arrow=False)
    if len(points) >= 2:
        add_line(slide, points[-2][0], points[-2][1], points[-1][0], points[-1][1], color=color, width=width, dash=dash, arrow=True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(PAPER)

    # Top banner.
    title = slide.shapes.add_textbox(Inches(0.35), Inches(0.14), Inches(9.6), Inches(0.38))
    set_text_frame(title, "MARLET Framework: Multi-Agent Retrieval Pipeline", 18, INK, True, PP_ALIGN.LEFT)
    subtitle = slide.shapes.add_textbox(Inches(0.38), Inches(0.54), Inches(9.6), Inches(0.26))
    set_text_frame(
        subtitle,
        "Coordination scopes the evidence need; retrieval is called only when it is worth spending.",
        9,
        MUTED,
        False,
        PP_ALIGN.LEFT,
    )
    rounded(slide, 10.25, 0.18, 2.65, 0.55, "Shared LLM backbone\nQwen3.5, budget 512, cap 900", WHITE, SLATE, size=8.5)

    # Main panels.
    group_band(slide, 0.38, 1.05, 1.72, 5.04, BLUE_LIGHT, BLUE, "1. Router")
    group_band(slide, 2.40, 1.05, 3.18, 5.04, GREEN_LIGHT, GREEN, "2. LLM specialists")
    group_band(slide, 5.88, 1.05, 1.74, 5.04, GOLD_LIGHT, GOLD, "3. Brief + gate")
    group_band(slide, 7.88, 1.05, 1.60, 5.04, TEAL_LIGHT, TEAL, "4. Retrieval")
    group_band(slide, 9.75, 1.05, 3.18, 5.04, PURPLE_LIGHT, PURPLE, "5. Prompt + output")

    # Router panel.
    inp = rounded(slide, 0.62, 1.42, 1.24, 0.54, "Input x\nq,o,c,h,r,v", WHITE, BLUE, 7.8)
    blob = rounded(slide, 0.62, 2.24, 1.24, 0.54, "Router blob\nq+c+r+h", WHITE, BLUE, 7.7)
    scores = rounded(slide, 0.62, 3.06, 1.24, 0.64, "Cue scores\ns_ev,s_cr,s_pl,s_ad", WHITE, BLUE, 6.9)
    route = rounded(slide, 0.62, 4.00, 1.24, 0.54, "RouteDecision\nR,G,roles", WHITE, BLUE, 7.5)
    add_line(slide, 1.24, 1.96, 1.24, 2.24, BLUE)
    add_line(slide, 1.24, 2.78, 1.24, 3.06, BLUE)
    add_line(slide, 1.24, 3.70, 1.24, 4.00, BLUE)

    # Specialists panel.
    base = rounded(slide, 2.72, 1.40, 2.54, 0.48, "Base AgentContext\nexample + route + budget", WHITE, GREEN, 7.4)
    planner = rounded(slide, 2.62, 2.30, 1.30, 0.62, "Planner LLM\nstrategy a\nqueries P", WHITE, GREEN, 6.8)
    diagnoser = rounded(slide, 4.00, 2.30, 1.30, 0.62, "Diagnoser LLM\nvisible state ell", WHITE, GREEN, 6.8)
    rubric = rounded(slide, 2.62, 3.42, 1.30, 0.62, "Criteria LLM\nrubric summary u", WHITE, GREEN, 6.8)
    tools = rounded(slide, 4.00, 3.42, 1.30, 0.62, "Tools T\nbounded local\nmax 4", WHITE, GREEN, 6.8)
    merge = rounded(slide, 2.80, 4.64, 2.38, 0.48, "Coordinator merge\nAgentResults -> response brief", GREEN_LIGHT, GREEN, 6.8)
    label(slide, 3.32, 1.96, 1.34, 0.18, "parallel async", 6.1, GREEN, True)
    add_polyline(slide, [(1.86, 4.27), (2.25, 4.27), (2.25, 1.64), (2.72, 1.64)], BLUE, 0.95)
    add_polyline(slide, [(3.99, 1.88), (3.99, 2.08), (3.27, 2.08), (3.27, 2.30)], GREEN, 0.85)
    add_polyline(slide, [(3.99, 1.88), (3.99, 2.08), (4.65, 2.08), (4.65, 2.30)], GREEN, 0.85)
    add_polyline(slide, [(3.99, 1.88), (3.99, 3.22), (3.27, 3.22), (3.27, 3.42)], GREEN, 0.85)
    add_polyline(slide, [(3.99, 1.88), (3.99, 3.22), (4.65, 3.22), (4.65, 3.42)], GREEN, 0.85, dash=True)
    add_line(slide, 3.27, 2.92, 3.27, 4.64, GREEN, 0.8, dash=True)
    add_line(slide, 4.65, 2.92, 4.65, 4.64, GREEN, 0.8, dash=True)
    add_line(slide, 3.27, 4.04, 3.27, 4.64, GREEN, 0.8, dash=True)
    add_line(slide, 4.65, 4.04, 4.65, 4.64, GREEN, 0.8, dash=True)

    # Brief + gate.
    brief = rounded(slide, 6.10, 2.02, 1.28, 0.70, "Response brief\np,ell,u,T", WHITE, GOLD, 7.1)
    gate = diamond(slide, 6.12, 3.35, 1.22, 0.86, "G=1?\ns_ev>=0.35\nor R=EG", GOLD_LIGHT, GOLD, 6.4)
    label(slide, 6.02, 4.34, 1.48, 0.20, "no: D is empty", 6.0, MUTED)
    add_line(slide, 5.18, 4.88, 6.10, 2.37, GREEN, 0.9, dash=True)
    add_line(slide, 6.74, 2.72, 6.74, 3.35, GOLD)

    # Retrieval branch.
    search = rounded(slide, 8.08, 1.78, 1.08, 0.52, "Search(I,P)\nplanner queries", WHITE, TEAL, 6.9)
    index = rounded(slide, 8.08, 2.72, 1.08, 0.52, "HybridIndex\nword+char+SVD", WHITE, TEAL, 6.7)
    pack = rounded(slide, 8.08, 3.66, 1.08, 0.56, "Dedup/Rerank\nMMR pack", WHITE, TEAL, 6.6)
    evidence = rounded(slide, 8.08, 4.64, 1.08, 0.46, "Evidence D\n<=4 chunks", TEAL_LIGHT, TEAL, 6.7)
    add_polyline(slide, [(7.34, 3.78), (7.76, 3.78), (7.76, 2.04), (8.08, 2.04)], TEAL, 1.0)
    add_line(slide, 8.62, 2.30, 8.62, 2.72, TEAL)
    add_line(slide, 8.62, 3.24, 8.62, 3.66, TEAL)
    add_line(slide, 8.62, 4.22, 8.62, 4.64, TEAL)

    # Generation panel.
    prompt = rounded(slide, 10.06, 1.42, 2.38, 0.72, "One assembled prompt\nq,o,h,p,ell,u,T,D,c", WHITE, PURPLE, 7.5)
    tutor = rounded(slide, 10.06, 2.62, 1.52, 0.60, "Tutor LLM\nwrites draft y0", WHITE, PURPLE, 7.5)
    fallback = diamond(slide, 11.76, 2.55, 0.74, 0.72, "Retry?\nD empty\ns_ev>=0.45", RED_LIGHT, RED, 5.5)
    critic = rounded(slide, 10.06, 4.36, 1.18, 0.56, "Critic LLM\none pass", WHITE, RED, 6.9)
    final = rounded(slide, 11.46, 4.30, 1.02, 0.68, "Final y\n+ trace", WHITE, RED, 6.9)
    add_polyline(slide, [(7.34, 3.78), (9.20, 3.78), (9.20, 1.78), (10.06, 1.78)], INK, 0.9)
    add_polyline(slide, [(9.16, 4.87), (9.48, 4.87), (9.48, 1.94), (10.06, 1.94)], TEAL, 0.9)
    add_line(slide, 11.25, 2.14, 11.25, 2.62, PURPLE)
    add_line(slide, 11.58, 2.92, 11.76, 2.91, RED, 0.9, dash=True)
    add_polyline(slide, [(12.13, 3.27), (12.13, 3.60), (9.60, 3.60), (9.60, 1.78), (8.08, 1.78)], RED, 0.85, dash=True)
    label(slide, 8.95, 3.34, 1.70, 0.18, "fallback uses Search(I,[q])", 5.8, RED)
    add_line(slide, 10.82, 3.22, 10.66, 4.36, INK, 0.85)
    add_line(slide, 11.24, 4.64, 11.46, 4.64, RED)

    # Clear notes.
    rounded(
        slide,
        2.70,
        5.36,
        2.80,
        0.44,
        "Agents communicate through structured records, not free-form chat.",
        WHITE,
        SLATE,
        6.8,
        bold=False,
    )
    rounded(slide, 10.02, 5.42, 2.58, 0.34, "solid = default path     dashed = conditional/fallback", WHITE, SLATE, 6.3, bold=False)

    # Footer citation to code path.
    label(
        slide,
        0.42,
        7.13,
        12.45,
        0.18,
        "Code-faithful flow from framework_codes/src/eduagentic/orchestration/pipelines.py::HybridFastPipeline and agents/{planner,diagnoser,rubric,tutor,critic}.py",
        5.8,
        MUTED,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)


if __name__ == "__main__":
    main()
